import os
from pptx import Presentation
from tkinter import Tk, filedialog, simpledialog, messagebox
from tkinterdnd2 import DND_FILES, TkinterDnD
from copy import deepcopy

def clone_slide(prs, slide):
    layout = prs.slide_layouts[slide.slide_layout.slide_layout_id]
    new_slide = prs.slides.add_slide(layout)

    for shape in slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(
            new_el, 'p:extLst'
        )

    for rel in slide.part.rels:
        if "notesSlide" not in rel:
            new_slide.part.rels.add_relationship(
                rel.reltype,
                rel._target,
                rel.rId
            )

def merge_presentations(files, output_file):
    merged = Presentation()

    # remove default slide
    while merged.slides:
        rId = merged.slides._sldIdLst[0].rId
        merged.part.drop_rel(rId)
        del merged.slides._sldIdLst[0]

    for file in files:
        prs = Presentation(file)
        for slide in prs.slides:
            clone_slide(merged, slide)

    merged.save(output_file)
    messagebox.showinfo("Success", f"Merged file created:\n{output_file}")

def select_folder():
    folder = filedialog.askdirectory(title="Select Folder Containing PPT Files")
    if not folder:
        return
    files = [os.path.join(folder, f) for f in os.listdir(folder) if f.endswith(".pptx")]
    process_files(files)

def process_files(files):
    if not files:
        messagebox.showerror("Error", "No PowerPoint files found.")
        return

    name = simpledialog.askstring("Output File", "Enter output file name:")
    if not name:
        return

    if not name.lower().endswith(".pptx"):
        name += ".pptx"

    output = os.path.join(os.path.dirname(files[0]), name)
    merge_presentations(files, output)

def drop(event):
    files = root.tk.splitlist(event.data)
    files = [f for f in files if f.endswith(".pptx")]
    process_files(files)

root = TkinterDnD.Tk()
root.title("PowerPoint Merger")
root.geometry("420x200")

label = messagebox.showinfo(
    "Instructions",
    "Drag & drop PPT files here\nOR\nCancel to select a folder"
)

root.drop_target_register(DND_FILES)
root.dnd_bind('<<Drop>>', drop)

if messagebox.askyesno("Folder Merge", "Do you want to merge an entire folder?"):
    select_folder()

root.mainloop()
